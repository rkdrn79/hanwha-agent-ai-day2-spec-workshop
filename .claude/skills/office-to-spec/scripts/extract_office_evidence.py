#!/usr/bin/env python3
"""Extract traceable raw evidence from PPTX and XLSX files."""

from __future__ import annotations

import argparse
from collections import Counter
import csv
import hashlib
import re
from pathlib import Path
from typing import Iterable

from openpyxl import load_workbook
from pptx import Presentation


def clean(value: object) -> str:
    if value is None:
        return ""
    return " ".join(str(value).replace("\n", " ").split())


def md_escape(value: object) -> str:
    return clean(value).replace("|", "\\|")


def safe_token(value: str) -> str:
    token = re.sub(r"[^A-Za-z0-9]+", "-", value).strip("-").upper()
    if token:
        return token
    return hashlib.sha1(value.encode("utf-8")).hexdigest()[:8].upper()


def locator_component(value: str) -> str:
    return value.replace("#", "_").replace("!", "_").replace(";", "_").replace("\n", " ")


def source_labels(paths: list[Path]) -> list[str]:
    """Return portable file labels, disambiguating repeated basenames."""
    counts = Counter(path.name for path in paths)
    labels: list[str] = []
    for path in paths:
        if counts[path.name] == 1:
            labels.append(path.name)
            continue
        digest = hashlib.sha1(str(path.resolve()).encode("utf-8")).hexdigest()[:8]
        labels.append(f"{path.stem}-{digest}{path.suffix}")
    return labels


def is_meaningful_ppt_evidence(text: str) -> bool:
    """Keep non-trivial source text without assuming a product or slide template."""
    normalized = text.strip()
    if not normalized or normalized in {"-", "–", "—", "•"}:
        return False
    if re.fullmatch(r"\d{1,4}", normalized):
        return False
    return len(normalized) >= 2


def iter_slide_blocks(slide) -> Iterable[tuple[str, str]]:
    block_no = 0
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = clean(shape.text)
            if text:
                block_no += 1
                yield f"B{block_no:02d}", text
        if getattr(shape, "has_table", False):
            block_no += 1
            table_id = f"T{block_no:02d}"
            for row_no, row in enumerate(shape.table.rows, start=1):
                values = [clean(cell.text) for cell in row.cells]
                if any(values):
                    yield f"{table_id}-R{row_no:02d}", " | ".join(values)


def extract_pptx(
    path: Path, source_label: str | None = None
) -> tuple[list[str], list[dict[str, str]]]:
    deck = Presentation(path)
    label = source_label or path.name
    file_token = safe_token(label)
    locator_file = locator_component(label)
    lines = [f"## PowerPoint: `{label}`", "", f"- 슬라이드 수: {len(deck.slides)}", ""]
    evidence: list[dict[str, str]] = []

    for slide_no, slide in enumerate(deck.slides, start=1):
        lines.extend([f"### PPTX:{locator_file}#S{slide_no}", ""])
        blocks = list(iter_slide_blocks(slide))
        if not blocks:
            lines.append("- 본문 텍스트 없음")
        for block_id, text in blocks:
            lines.append(f"- `{block_id}` {md_escape(text)}")
        for block_id, text in blocks:
            if not is_meaningful_ppt_evidence(text):
                continue
            evidence.append(
                {
                    "evidence_id": f"PPTX-{file_token}-S{slide_no:02d}-{block_id}",
                    "source_file": label,
                    "locator": f"PPTX:{locator_file}#S{slide_no}:{block_id}",
                    "evidence_kind": "slide_table_row" if block_id.startswith("T") else "slide_text",
                    "raw_text": text,
                    "classification": "",
                    "normalized_claim": "",
                    "requirement_id": "",
                    "notes": f"parent=PPTX:{locator_file}#S{slide_no}",
                }
            )

        notes_text = ""
        try:
            notes_text = clean(slide.notes_slide.notes_text_frame.text)
        except (AttributeError, ValueError):
            notes_text = ""
        if notes_text:
            lines.append(f"- `NOTES` {md_escape(notes_text)}")
            evidence.append(
                {
                    "evidence_id": f"PPTX-{file_token}-S{slide_no:02d}-NOTES",
                    "source_file": label,
                    "locator": f"PPTX:{locator_file}#S{slide_no}:NOTES",
                    "evidence_kind": "speaker_notes",
                    "raw_text": notes_text,
                    "classification": "",
                    "normalized_claim": "",
                    "requirement_id": "",
                    "notes": "",
                }
            )
        lines.append("")

    return lines, evidence


def extract_xlsx(
    path: Path, source_label: str | None = None
) -> tuple[list[str], list[dict[str, str]]]:
    wb = load_workbook(path, data_only=False)
    label = source_label or path.name
    file_token = safe_token(label)
    locator_file = locator_component(label)
    lines = [f"## Excel: `{label}`", "", f"- 시트 수: {len(wb.worksheets)}", ""]
    evidence: list[dict[str, str]] = []

    for ws in wb.worksheets:
        sheet_token = safe_token(ws.title)
        locator_sheet = locator_component(ws.title)
        lines.extend(
            [
                f"### XLSX:{locator_file}#{locator_sheet}",
                "",
                f"- 상태: `{ws.sheet_state}`",
                f"- 사용 범위: `{ws.calculate_dimension()}`",
            ]
        )
        if ws.merged_cells.ranges:
            merged = ", ".join(str(item) for item in ws.merged_cells.ranges)
            lines.append(f"- 병합 범위: `{merged}`")
        if ws.data_validations.count:
            validations = []
            for validation in ws.data_validations.dataValidation:
                validations.append(
                    f"{validation.sqref} ({validation.type}: {validation.formula1 or ''})"
                )
            lines.append(f"- 데이터 검증: `{'; '.join(validations)}`")
        lines.append("")

        for row_no in range(1, ws.max_row + 1):
            cells = [cell for cell in ws[row_no] if cell.value is not None]
            if not cells:
                continue
            start = min(cell.column for cell in cells)
            end = max(cell.column for cell in cells)
            start_coord = ws.cell(row_no, start).coordinate
            end_coord = ws.cell(row_no, end).coordinate
            locator = f"XLSX:{locator_file}#{locator_sheet}!{start_coord}:{end_coord}"
            parts = []
            comments = []
            for cell in cells:
                tags = []
                if cell.data_type == "f":
                    tags.append("(formula)")
                validations = [
                    validation
                    for validation in ws.data_validations.dataValidation
                    if cell.coordinate in validation.sqref
                ]
                if validations:
                    tags.append("(validation)")
                    details = "; ".join(
                        f"{validation.type}:{clean(validation.formula1)}"
                        for validation in validations
                    )
                    tags.append(f"[{details}]")
                suffix = f" {' '.join(tags)}" if tags else ""
                parts.append(f"{cell.coordinate}={clean(cell.value)}{suffix}")
                if cell.comment and clean(cell.comment.text):
                    comments.append(f"{cell.coordinate} comment={clean(cell.comment.text)}")
            raw_text = " | ".join(parts)
            if comments:
                raw_text += " | " + " | ".join(comments)
            lines.append(f"- `{locator}` {md_escape(raw_text)}")
            evidence.append(
                {
                    "evidence_id": f"XLSX-{file_token}-{sheet_token}-R{row_no:03d}",
                    "source_file": label,
                    "locator": locator,
                    "evidence_kind": "excel_row",
                    "raw_text": raw_text,
                    "classification": "",
                    "normalized_claim": "",
                    "requirement_id": "",
                    "notes": "",
                }
            )
        lines.append("")

    return lines, evidence


def write_outputs(
    output_dir: Path,
    source_paths: list[Path],
    inventory_lines: list[str],
    evidence: list[dict[str, str]],
) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    inventory = [
        "# Office source inventory",
        "",
        "> 자동 추출 결과입니다. 슬라이드 렌더링과 실제 셀·수식 확인을 함께 수행하세요.",
        "",
        "## 입력 파일",
        "",
        *(f"- `{path}`" for path in source_paths),
        "",
        *inventory_lines,
    ]
    (output_dir / "01_source_inventory.md").write_text(
        "\n".join(inventory).rstrip() + "\n", encoding="utf-8"
    )

    fields = [
        "evidence_id",
        "source_file",
        "locator",
        "evidence_kind",
        "raw_text",
        "classification",
        "normalized_claim",
        "requirement_id",
        "notes",
    ]
    with (output_dir / "02_evidence_ledger.csv").open(
        "w", encoding="utf-8-sig", newline=""
    ) as handle:
        writer = csv.DictWriter(handle, fieldnames=fields)
        writer.writeheader()
        writer.writerows(evidence)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--pptx", type=Path, action="append", default=[])
    parser.add_argument("--xlsx", type=Path, action="append", default=[])
    parser.add_argument("--output-dir", type=Path, required=True)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    source_paths = [*args.pptx, *args.xlsx]
    if not source_paths:
        raise SystemExit("하나 이상의 --pptx 또는 --xlsx 입력이 필요합니다.")
    resolved_paths = [path.resolve() for path in source_paths]
    if len(set(resolved_paths)) != len(resolved_paths):
        raise SystemExit("같은 입력 파일이 두 번 이상 지정되었습니다.")
    for path in source_paths:
        if not path.is_file():
            raise SystemExit(f"입력 파일을 찾을 수 없습니다: {path}")

    inventory_lines: list[str] = []
    evidence: list[dict[str, str]] = []
    labels = source_labels(source_paths)
    pptx_labels = labels[: len(args.pptx)]
    xlsx_labels = labels[len(args.pptx) :]
    for path, label in zip(args.pptx, pptx_labels, strict=True):
        lines, rows = extract_pptx(path, label)
        inventory_lines.extend(lines)
        evidence.extend(rows)
    for path, label in zip(args.xlsx, xlsx_labels, strict=True):
        lines, rows = extract_xlsx(path, label)
        inventory_lines.extend(lines)
        evidence.extend(rows)
    write_outputs(
        args.output_dir,
        source_paths,
        inventory_lines,
        evidence,
    )
    print(f"created: {args.output_dir / '01_source_inventory.md'}")
    print(f"created: {args.output_dir / '02_evidence_ledger.csv'}")
    print(f"evidence rows: {len(evidence)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

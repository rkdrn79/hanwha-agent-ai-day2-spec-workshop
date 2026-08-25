#!/usr/bin/env python3
"""Verify the committed Spec ↔ Office workshop artifacts."""

from __future__ import annotations

import json
import re
import subprocess
import sys
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
ERROR_VALUES = {"#REF!", "#DIV/0!", "#VALUE!", "#NAME?", "#N/A"}
checks: list[str] = []


def require(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)
    checks.append(message)


def verify_files() -> None:
    required = [
        "presentation/output/ips_wbs_solution_brief.pptx",
        "presentation/output/contact_sheet.png",
        "reverse-spec/input/roomflow_screen_definition.pptx",
        "reverse-spec/input/roomflow_function_definition.xlsx",
        "reverse-spec/previews/roomflow_ppt_contact.jpg",
        "reverse-spec/previews/roomflow_요약.png",
        "reverse-spec/expected/SPEC.md",
        "reverse-spec/expected/TRACEABILITY.csv",
        "reverse-spec/expected/OPEN_QUESTIONS.md",
        ".claude/skills/office-to-spec/SKILL.md",
        ".claude/skills/pptx/SKILL.md",
        ".claude/skills/xlsx/SKILL.md",
    ]
    missing = [item for item in required if not (ROOT / item).is_file()]
    require(not missing, "필수 실습 파일이 모두 존재함")


def verify_presentations() -> None:
    spec_to_ppt = Presentation(ROOT / "presentation/output/ips_wbs_solution_brief.pptx")
    require(len(spec_to_ppt.slides) == 8, "Spec → PPT 완성본이 8장임")

    source_ppt = Presentation(ROOT / "reverse-spec/input/roomflow_screen_definition.pptx")
    require(len(source_ppt.slides) == 8, "Office → Spec 입력 PPT가 8장임")
    notes_count = 0
    for slide in source_ppt.slides:
        try:
            if slide.notes_slide.notes_text_frame.text.strip():
                notes_count += 1
        except (AttributeError, ValueError):
            pass
    require(notes_count == 8, "입력 PPT 8장 모두 발표자 노트가 있음")


def verify_workbook() -> None:
    path = ROOT / "reverse-spec/input/roomflow_function_definition.xlsx"
    formulas = load_workbook(path, data_only=False)
    values = load_workbook(path, data_only=True)
    expected_sheets = [
        "요약",
        "화면목록",
        "화면항목정의",
        "버튼이벤트",
        "상태전이",
        "메시지·권한",
        "검토메모",
    ]
    require(formulas.sheetnames == expected_sheets, "입력 Excel의 7개 시트와 순서가 정확함")

    formula_cells = [
        cell.coordinate
        for sheet in formulas.worksheets
        for row in sheet.iter_rows()
        for cell in row
        if cell.data_type == "f"
    ]
    require(len(formula_cells) == 7, "입력 Excel에 자동 집계 수식 7개가 있음")

    errors = [
        sheet.title + "!" + cell.coordinate
        for sheet in values.worksheets
        for row in sheet.iter_rows()
        for cell in row
        if cell.value in ERROR_VALUES
    ]
    require(not errors, "입력 Excel의 계산 결과에 수식 오류가 없음")
    require(
        [values["요약"][cell].value for cell in ("E5", "E6", "E7", "E8", "E9", "E10", "E11")]
        == [4, 12, 9, 7, 5, 4, 6],
        "Excel 요약 수식 결과가 4/12/9/7/5/4/6임",
    )
    require(
        values["화면목록"]["E5"].value == "오늘부터 30일"
        and values["화면항목정의"]["E14"].value == "N"
        and values["버튼이벤트"]["D8"].value == "APPROVED이며 시작 4시간 전까지",
        "Excel의 의도된 충돌 근거가 보존됨",
    )


def verify_expected_spec() -> None:
    command = [
        sys.executable,
        str(ROOT / ".claude/skills/office-to-spec/scripts/validate_spec.py"),
        "--spec",
        str(ROOT / "reverse-spec/expected/SPEC.md"),
        "--traceability",
        str(ROOT / "reverse-spec/expected/TRACEABILITY.csv"),
        "--questions",
        str(ROOT / "reverse-spec/expected/OPEN_QUESTIONS.md"),
    ]
    result = subprocess.run(command, check=False, capture_output=True, text=True)
    if result.returncode != 0:
        raise AssertionError("정답 SPEC 검증 실패: " + result.stdout + result.stderr)
    payload = json.loads(result.stdout)
    require(payload["status"] == "passed", "정답 SPEC 추적성 검증이 passed임")
    require(payload["spec_requirement_count"] == 21, "정답 SPEC의 FR/AC가 21개임")


def verify_readme_prompts() -> None:
    readme = (ROOT / "README.md").read_text(encoding="utf-8")
    required_markers = [
        "# Part 1. Spec → PPT",
        "# Part 2. PPT + Excel → Spec",
        "## Prompt 1 — 원본 조사만 하기",
        "## Prompt 2 — 증거를 분류하고 질문 만들기",
        "## Prompt 3 — 확정 증거만 SPEC으로 만들기",
        "## Prompt 4 — 자동 검증과 red-team",
    ]
    require(all(marker in readme for marker in required_markers), "README에 순서형 Claude 프롬프트가 있음")


def verify_markdown_links() -> None:
    missing: list[str] = []
    pattern = re.compile(r"\[[^\]]+\]\(([^)]+)\)")
    for markdown_path in (ROOT / "README.md", ROOT / "presentation/README.md", ROOT / "reverse-spec/README.md"):
        text = markdown_path.read_text(encoding="utf-8")
        for raw_target in pattern.findall(text):
            target = raw_target.strip().split("#", 1)[0]
            if not target or target.startswith(("http://", "https://", "mailto:")):
                continue
            if not (markdown_path.parent / target).exists():
                missing.append(str(markdown_path.relative_to(ROOT)) + " -> " + target)
    require(not missing, "주요 README의 로컬 링크가 모두 유효함")


def main() -> int:
    try:
        verify_files()
        verify_presentations()
        verify_workbook()
        verify_expected_spec()
        verify_readme_prompts()
        verify_markdown_links()
    except AssertionError as error:
        print(json.dumps({"status": "failed", "error": str(error), "checks": checks}, ensure_ascii=False, indent=2))
        return 1

    print(json.dumps({"status": "passed", "check_count": len(checks), "checks": checks}, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
